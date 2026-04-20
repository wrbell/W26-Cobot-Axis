#!/usr/bin/env python3
"""Render reports/turn-in/presentation/presentation.md -> presentation.pptx via python-pptx.

Convention:
  # H1              -> new slide (slide title; subtitle below it via ## if present)
  ## H2             -> subtitle / section header on the slide
  ### H3            -> bold small-cap header within content
  - bullet          -> bullet list item
  | .. | .. |       -> table (inserted beneath bullets)
  <!-- NOTES: ... -->  -> speaker notes (multiline OK; trailing -->)
  <!-- other -->    -> stripped

A standalone `---` is treated as a slide break for safety even though H1 is
the primary separator.
"""
from __future__ import annotations

import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

HERE = Path(__file__).parent
SRC = HERE / "presentation.md"
OUT = HERE / "presentation.pptx"

NOTES_RE = re.compile(r"<!--\s*NOTES:\s*(.*?)\s*-->", re.DOTALL)
PLAIN_COMMENT = re.compile(r"<!--(?!\s*NOTES:).*?-->", re.DOTALL)


class Slide:
    def __init__(self) -> None:
        self.title: str = ""
        self.subtitle: str = ""
        self.bullets: list[str] = []
        self.table_rows: list[list[str]] = []
        self.notes: str = ""
        self.body_lines: list[str] = []  # free-form paragraphs (non-bullet)


def split_slides(text: str) -> list[Slide]:
    # Extract NOTES first so we can attach them to the slide they belong to.
    slides: list[Slide] = []
    current: Slide | None = None
    lines = text.splitlines()
    i = 0
    in_table = False
    table_block: list[str] = []

    def flush_table():
        nonlocal table_block, in_table
        if table_block and current is not None:
            rows = []
            for line in table_block:
                if set(line.replace("|", "").replace(":", "").replace("-", "").strip()) == set():
                    continue
                cells = [c.strip() for c in line.strip().strip("|").split("|")]
                rows.append(cells)
            if rows:
                current.table_rows = rows
        table_block = []
        in_table = False

    while i < len(lines):
        line = lines[i]

        # Speaker notes are multi-line comment blocks
        if "<!--" in line and "NOTES:" in line:
            # gather until -->
            buf = [line]
            while "-->" not in buf[-1] and i + 1 < len(lines):
                i += 1
                buf.append(lines[i])
            full = "\n".join(buf)
            m = NOTES_RE.search(full)
            if m and current is not None:
                current.notes = (current.notes + "\n\n" + m.group(1).strip()).strip()
            i += 1
            continue

        # Other HTML comments -> strip silently
        if line.strip().startswith("<!--") and line.strip().endswith("-->"):
            i += 1
            continue

        # Slide break
        if line.strip() == "---":
            flush_table()
            current = None
            i += 1
            continue

        # H1 = new slide with title
        if line.startswith("# "):
            flush_table()
            current = Slide()
            current.title = line[2:].strip()
            slides.append(current)
            i += 1
            continue

        # H2 = subtitle (on current slide)
        if line.startswith("## "):
            if current is None:
                current = Slide()
                slides.append(current)
            # First ## becomes subtitle; additional become bold headers in body
            if not current.subtitle:
                current.subtitle = line[3:].strip()
            else:
                current.body_lines.append(f"**{line[3:].strip()}**")
            i += 1
            continue

        if line.startswith("### "):
            if current is not None:
                current.body_lines.append(f"**{line[4:].strip()}**")
            i += 1
            continue

        # Table rows
        if line.strip().startswith("|") and line.strip().endswith("|"):
            in_table = True
            table_block.append(line)
            i += 1
            continue
        elif in_table:
            flush_table()

        # Bullet
        stripped = line.lstrip()
        if stripped.startswith(("- ", "* ", "+ ")):
            if current is not None:
                current.bullets.append(stripped[2:].strip())
            i += 1
            continue

        # Nested bullets (4+ space indent)
        if line.startswith("    ") and line.lstrip().startswith(("- ", "* ", "+ ")):
            if current is not None and current.bullets:
                current.bullets[-1] = current.bullets[-1] + "\n    " + line.lstrip()
            i += 1
            continue

        # Blank line
        if not line.strip():
            i += 1
            continue

        # Plain paragraph line
        if current is not None:
            current.body_lines.append(line.strip())
        i += 1

    flush_table()
    return slides


def apply_inline(run_paragraph, text: str) -> None:
    """Minimal inline formatter for **bold** and *italic* within a pptx paragraph."""
    pattern = re.compile(r"(\*\*[^*]+\*\*|\*[^*]+\*)")
    pos = 0
    for m in pattern.finditer(text):
        if m.start() > pos:
            run = run_paragraph.add_run()
            run.text = text[pos:m.start()]
        token = m.group(0)
        run = run_paragraph.add_run()
        if token.startswith("**"):
            run.text = token[2:-2]
            run.font.bold = True
        else:
            run.text = token[1:-1]
            run.font.italic = True
        pos = m.end()
    if pos < len(text):
        run = run_paragraph.add_run()
        run.text = text[pos:]


def build_slide(prs: Presentation, slide_data: Slide) -> None:
    blank = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(blank)

    # Title
    if slide.shapes.title is not None:
        slide.shapes.title.text = slide_data.title or "(Untitled)"
        for p in slide.shapes.title.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(32)
                r.font.bold = True

    # Body text box
    left = Inches(0.5)
    top = Inches(1.6)
    width = Inches(9.0)
    height = Inches(5.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    first = True

    def new_para(text: str = "", bullet: bool = False, size: int = 18):
        nonlocal first
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.level = 1 if bullet else 0
        if text:
            apply_inline(p, text)
            for r in p.runs:
                r.font.size = Pt(size)
        return p

    # Subtitle
    if slide_data.subtitle:
        p = new_para(slide_data.subtitle, size=20)
        for r in p.runs:
            r.font.italic = True

    # Free-form body lines (before bullets)
    for line in slide_data.body_lines:
        new_para(line, size=18)

    # Bullets
    for b in slide_data.bullets:
        new_para(b, bullet=True, size=18)

    # Table (inserted below body if present)
    if slide_data.table_rows:
        rows = slide_data.table_rows
        tbl_top = Inches(2.8) if slide_data.bullets or slide_data.body_lines else Inches(1.8)
        table_shape = slide.shapes.add_table(
            rows=len(rows), cols=len(rows[0]),
            left=Inches(0.5), top=tbl_top,
            width=Inches(9.0), height=Inches(0.4 * len(rows)),
        ).table
        for r_idx, row in enumerate(rows):
            for c_idx, cell_text in enumerate(row):
                cell = table_shape.cell(r_idx, c_idx)
                cell.text = ""
                p = cell.text_frame.paragraphs[0]
                apply_inline(p, cell_text)
                for run in p.runs:
                    run.font.size = Pt(12)
                    if r_idx == 0:
                        run.font.bold = True

    # Speaker notes
    if slide_data.notes:
        notes_tf = slide.notes_slide.notes_text_frame
        notes_tf.text = slide_data.notes


def main() -> int:
    if not SRC.exists():
        print(f"ERROR: {SRC} not found", file=sys.stderr)
        return 1

    text = SRC.read_text(encoding="utf-8")
    # NB: do NOT strip NOTES comments here — split_slides handles them itself
    slides = split_slides(text)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    for slide in slides:
        build_slide(prs, slide)

    prs.save(OUT)
    print(f"Wrote {OUT.name} ({OUT.stat().st_size} bytes, {len(slides)} slides)")
    for idx, s in enumerate(slides, start=1):
        has_notes = "N" if s.notes else " "
        has_tbl = "T" if s.table_rows else " "
        print(f"  {idx:>2}. [{has_notes}{has_tbl}] {s.title}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
